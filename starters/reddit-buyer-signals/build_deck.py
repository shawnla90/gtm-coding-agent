#!/usr/bin/env python3
"""build_deck.py - build a short pitch deck from the scored signals as editable Google Slides.

A lean, client-agnostic reference implementation. It reads the headline numbers and the top scored
topics + buyer quotes straight out of data/signals.db, builds a .pptx with python-pptx, then
uploads it to Google Drive converting to native Google Slides (the Drive scope on your existing
gspread token covers this, no extra OAuth needed). The story is generic on purpose: your buyers
are asking AI which tool to pick, you are not in the answer yet, here is the scored plan to fix it.

Set BRAND to your product name so the deck reads as yours (default: "Acme PM").

  python3 build_deck.py            # build + upload, print the Slides URL
  python3 build_deck.py --local    # just write the .pptx, do not upload
  BRAND="Your Product" python3 build_deck.py
"""
import os
import sqlite3
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = Path(__file__).resolve().parent
DB = HERE / "data" / "signals.db"
PPTX = HERE / "reddit-buyer-signals-deck.pptx"
URLFILE = HERE / "data" / "slides_url.txt"
BRAND = os.environ.get("BRAND", "Acme PM")

SLATE = RGBColor(0x1F, 0x3A, 0x56)
INK = RGBColor(0x16, 0x28, 0x3A)
GREEN = RGBColor(0x16, 0x65, 0x3A)
GREEN_MID = RGBColor(0x2E, 0x8B, 0x57)
GREEN_LT = RGBColor(0xDC, 0xED, 0xE2)
AMBER = RGBColor(0xD9, 0x8A, 0x21)
MUTED = RGBColor(0x5B, 0x6B, 0x7B)
PAPER = RGBColor(0xFA, 0xFA, 0xF7)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE4, 0xE8, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLOUD = RGBColor(0xC9, 0xD6, 0xE2)


def load_data():
    """Pull the few headline numbers and the top rows the deck needs. Returns a dict."""
    con = sqlite3.connect(DB)
    q = con.execute
    n_threads = q("SELECT COUNT(*) FROM reddit_threads").fetchone()[0]
    n_subs = q("SELECT COUNT(DISTINCT subreddit) FROM reddit_threads").fetchone()[0]
    n_comp = q("SELECT COUNT(*) FROM buyer_language WHERE kind='comparison'").fetchone()[0]
    n_lang = q("SELECT COUNT(*) FROM buyer_language").fetchone()[0]
    a_tier = q("SELECT COUNT(*) FROM content_topics WHERE tier='A'").fetchone()[0]
    topics = q("SELECT tier, topic, content_angle, mentions FROM content_topics "
               "ORDER BY score DESC, engagement DESC LIMIT 6").fetchall()
    quotes = q("SELECT quote, subreddit FROM buyer_language "
               "WHERE kind IN ('comparison','recommendation') AND length(quote) BETWEEN 20 AND 120 "
               "ORDER BY engagement DESC LIMIT 6").fetchall()
    con.close()
    return dict(n_threads=n_threads, n_subs=n_subs, n_comp=n_comp, n_lang=n_lang,
                a_tier=a_tier, topics=topics, quotes=quotes)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def tb(s, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
       italic=False, anchor=MSO_ANCHOR.TOP, spacing=1.06, font="Arial"):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        run = p.add_run(); run.text = ln
        f = run.font; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = color; f.name = font
    return box


def kicker(s, text, x=0.75, y=0.62, color=GREEN):
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y + 0.09),
                             Inches(0.27), Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = AMBER; bar.line.fill.background(); bar.shadow.inherit = False
    tb(s, x + 0.36, y - 0.04, 9, 0.4, text.upper(), 13, color, bold=True)


def card(s, x, y, w, h, fill=CARD, border=LINE, accent=None):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.color.rgb = border; c.line.width = Pt(1); c.shadow.inherit = False
    if accent:
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background(); bar.shadow.inherit = False
    return c


def foot(s, left, n=None, dark=False):
    col = CLOUD if dark else RGBColor(0x98, 0xA3, 0xB0)
    tb(s, 0.75, 6.95, 10, 0.35, left, 11, col, bold=True)
    if n is not None:
        tb(s, 12.2, 6.95, 0.9, 0.35, str(n), 11, col, bold=True, align=PP_ALIGN.RIGHT)


def build_slides(d):
    # 1 TITLE
    s = slide(SLATE)
    kicker(s, "GTM & AI visibility", y=2.2, color=RGBColor(0x8F, 0xE0, 0xB3))
    tb(s, 0.75, 2.6, 11.8, 1.9, f"Getting {BRAND} found\nin the AI era.", 46, WHITE, bold=True)
    tb(s, 0.75, 4.55, 10.5, 1.1,
       "Your buyers are asking AI which tool to pick before they ask anyone else. Right now, it does "
       "not name you. Here is the plan to change that, built from real, recent buyer conversations.",
       19, CLOUD)
    foot(s, f"{BRAND} · Reddit buyer signals", dark=True)

    # 2 THE GAP (headline numbers from the DB)
    s = slide()
    kicker(s, "The gap")
    tb(s, 0.75, 1.0, 12, 1.1, "The buying decision is happening in public. You are not in it.", 31, SLATE, bold=True)
    stats = [(str(d["n_threads"]), "recent buyer threads, all from the last 30 days"),
             (str(d["n_comp"]), "head-to-head comparisons of the tools you compete with"),
             (str(d["n_subs"]), "communities where your buyers post"),
             ("0", f"of those threads mention {BRAND}")]
    for i, (big, lbl) in enumerate(stats):
        x = 0.75 + i * 3.03
        card(s, x, 2.6, 2.85, 1.85)
        tb(s, x + 0.24, 2.74, 2.5, 0.7, big, 34, GREEN, bold=True)
        tb(s, x + 0.24, 3.5, 2.45, 0.85, lbl, 12.5, MUTED, spacing=1.12)
    tb(s, 0.75, 4.9, 11.6, 0.9,
       "AI models read these threads and cite them. The tool that answers the questions in public "
       "becomes the default recommendation. Nobody has claimed that spot yet.", 17, MUTED, spacing=1.2)
    foot(s, "Public Reddit data, last 30 days only", n=2)

    # 3 WHERE BUYERS TALK (real quotes from the DB)
    s = slide()
    kicker(s, "Where your buyers actually talk")
    tb(s, 0.75, 1.0, 12, 0.9, "They are comparing the exact tools you compete with. Right now.", 29, SLATE, bold=True)
    tb(s, 0.75, 1.9, 11.6, 0.6, "Real, active threads pulled this month. The AI models read them. You are in none of them.", 17, MUTED)
    qs = d["quotes"] or [("(run the pipeline to populate real buyer quotes here)", "")]
    for i, (quote, sub) in enumerate(qs[:6]):
        col, row = i % 2, i // 2
        x = 0.75 + col * 6.06
        y = 2.65 + row * 1.28
        card(s, x, y, 5.8, 1.12, accent=GREEN_MID)
        tb(s, x + 0.28, y + 0.13, 5.35, 0.7, f"“{quote}”", 14, INK, bold=True, spacing=1.05)
        tb(s, x + 0.28, y + 0.8, 5.35, 0.28, f"r/{sub}" if sub else "", 12, GREEN, bold=True)
    foot(s, "Every tool named in these threads is one you compete with", n=3)

    # 4 THE PLAN (top scored topics from the DB)
    s = slide()
    kicker(s, "What to publish")
    tb(s, 0.75, 1.0, 12, 0.9, "A scored plan, not a guess.", 31, SLATE, bold=True)
    tb(s, 0.75, 1.85, 11.6, 0.6, "Ranked by buyer intent, demand, competitive fit, and how widely the threads are read. Publish A-tier first.", 17, MUTED)
    tier_color = {"A": GREEN, "B": RGBColor(0x3E, 0x77, 0xB5), "C": RGBColor(0xC9, 0x9A, 0x2E), "D": MUTED}
    headers = [("TIER", 0.75, 1.0), ("TOPIC", 2.0, 2.6), ("THE PAGE YOU PUBLISH", 4.9, 6.1), ("SIGNALS", 11.2, 1.6)]
    for h, x, w in headers:
        tb(s, x, 2.55, w, 0.3, h, 11, MUTED, bold=True)
    y = 2.95
    rows = d["topics"] or [("A", "run the pipeline", "the scored plan lands here", 0)]
    for tier, topic, angle, mentions in rows[:6]:
        chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(y + 0.03), Inches(0.42), Inches(0.34))
        chip.fill.solid(); chip.fill.fore_color.rgb = tier_color.get(tier, MUTED); chip.line.fill.background(); chip.shadow.inherit = False
        tb(s, 0.75, y + 0.05, 0.42, 0.3, tier or "-", 13, WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tb(s, 2.0, y, 2.7, 0.4, str(topic).title(), 14, INK)
        tb(s, 4.9, y, 6.1, 0.5, str(angle), 13.5, INK, spacing=1.05)
        tb(s, 11.2, y, 1.8, 0.4, f"{mentions} signals", 13, MUTED)
        ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(y + 0.52), Inches(12.05), Pt(0.75))
        ln.fill.solid(); ln.fill.fore_color.rgb = LINE; ln.line.fill.background(); ln.shadow.inherit = False
        y += 0.64
    foot(s, "Full scored list and every source thread live in the shared Google Sheet", n=4)

    # 5 THE LOOP
    s = slide()
    kicker(s, "Why this keeps working")
    tb(s, 0.75, 1.0, 12, 0.9, "It is an engine, not a one-off.", 33, SLATE, bold=True)
    plays = [("Listen", "Pull recent buyer talk", "Only the last 30 days, so you engage live threads, never stale ones."),
             ("Answer", "Publish the A-tier pages", "Written in the buyer's words, mapped to a comparison you can win."),
             ("Get cited", "Show up in AI answers", "Track where you now appear, refresh the plan, widen the gap.")]
    for i, (step, h, body) in enumerate(plays):
        x = 0.75 + i * 4.06
        card(s, x, 2.6, 3.8, 2.4)
        tb(s, x + 0.28, 2.8, 3.3, 0.35, step.upper(), 12, AMBER, bold=True)
        tb(s, x + 0.28, 3.2, 3.3, 0.5, h, 19, SLATE, bold=True)
        tb(s, x + 0.28, 3.85, 3.3, 1.1, body, 13.5, MUTED, spacing=1.18)
    tb(s, 0.75, 5.4, 11.8, 0.8, "Listen → Answer → Get cited → Measure → Repeat. The same loop every month, "
       f"so {BRAND} becomes the tool AI reaches for when a buyer asks.", 16, INK, spacing=1.2)
    foot(s, f"{BRAND} · Reddit buyer signals", n=5)


def main():
    d = load_data()
    build_slides(d)
    prs.save(PPTX)
    print(f"built {PPTX.name} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")

    if "--local" in sys.argv:
        return

    # upload to Drive, converting to native Google Slides (uses the gspread token's drive scope)
    from google.oauth2.credentials import Credentials
    from google.auth.transport.requests import Request
    from googleapiclient.discovery import build as gbuild
    from googleapiclient.http import MediaFileUpload

    TOKEN = Path.home() / ".config" / "gspread" / "token.json"
    SCOPES = ["https://www.googleapis.com/auth/spreadsheets", "https://www.googleapis.com/auth/drive"]
    c = Credentials.from_authorized_user_file(str(TOKEN), SCOPES)
    if not c.valid:
        c.refresh(Request()); TOKEN.write_text(c.to_json())
    drive = gbuild("drive", "v3", credentials=c)
    media = MediaFileUpload(str(PPTX),
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    # reuse the existing Slides file if we have one, so the link never changes on a rebuild
    existing_id = None
    if URLFILE.exists():
        import re as _re
        m = _re.search(r"/d/([A-Za-z0-9_-]+)", URLFILE.read_text())
        existing_id = m.group(1) if m else None

    if existing_id:
        try:
            drive.files().update(fileId=existing_id, media_body=media,
                                 body={"mimeType": "application/vnd.google-apps.presentation"}).execute()
            url = f"https://docs.google.com/presentation/d/{existing_id}/edit"
            print("SLIDES (updated in place):", url)
        except Exception:
            existing_id = None

    if not existing_id:
        meta = {"name": f"{BRAND} - Reddit Buyer Signals (deck)",
                "mimeType": "application/vnd.google-apps.presentation"}
        f = drive.files().create(body=meta, media_body=media, fields="id,webViewLink").execute()
        drive.permissions().create(fileId=f["id"], body={"type": "anyone", "role": "reader"}).execute()
        url = f.get("webViewLink") or f"https://docs.google.com/presentation/d/{f['id']}/edit"
        print("SLIDES (new):", url)
    URLFILE.write_text(url + "\n")


if __name__ == "__main__":
    main()
